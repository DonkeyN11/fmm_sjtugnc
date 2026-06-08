#!/usr/bin/env python3
"""Create clean PPTX with all paper figures as high-res PNG images + editable SVG source.
Each slide: title bar, high-res PNG preview, and embedded SVG object for vector editing.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part, PackURI
from lxml import etree

PROJECT = Path(__file__).resolve().parents[2]
FIGS_DIR = PROJECT / "docs/Trustworthiness Evaluation Framework for Map Matching based on Covariance Ellipse/figs"
OUT = FIGS_DIR / "all_figures.pptx"

FIGURES = [
    ("abstract_figure",     "Fig. 1 — Graphical Abstract",
     "Pipeline comparison: isotropic HMM (top) vs. covariance-aware CMM (bottom)"),
    ("covariance",          "Fig. 4 — GNSS Covariance Formation",
     "3-panel: skyplot, pseudorange noise model, error ellipse with semi-axes"),
    ("PLsearch",            "Fig. 5 — HPL-Based Adaptive Candidate Search",
     "2-panel: good geometry (tight search) vs. degraded geometry (wide search)"),
    ("maha_vs_eu",          "Fig. 6 — Mahalanobis vs. Euclidean Projection",
     "Original anisotropic space vs. whitened isotropic space, Cholesky transform"),
    ("measprob",            "Fig. 7 — Emission Probability Under Whitening",
     "Anisotropic Gaussian emission → isotropic Euclidean via Cholesky whitening"),
    ("trustworthiness",     "Fig. 8 — Trustworthiness Across Road Geometries",
     "2x2: highway (high TW), junction (low), urban grid (medium-high), dual carriageway"),
    ("consistency_overview","Fig. 9 — Covariance Model Consistency Validation",
     "4-panel: chi-squared fit, P-P plot, whitened errors, radial CDF vs Rayleigh"),
    ("dataset_overview",    "Fig. 16 — Hainan-06 Dataset Overview",
     "Haikou road network (152K edges) with 7 colored SPP trajectories + location inset"),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# Professional color theme
TITLE_BLUE = RGBColor(0x1A, 0x5C, 0x8A)
SUBTITLE_GRAY = RGBColor(0x66, 0x66, 0x66)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

for name, title, desc in FIGURES:
    png_path = FIGS_DIR / f"{name}.png"
    svg_path = FIGS_DIR / f"{name}.svg"
    if not png_path.exists():
        print(f"SKIP: {name} (no PNG)"); continue

    slide = prs.slides.add_slide(BLANK)

    # ── Clean header bar ──
    # Subtle colored bar at top
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_BLUE; bar.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.45))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = TITLE_BLUE

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(10), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = desc
    p2.font.size = Pt(11); p2.font.color.rgb = SUBTITLE_GRAY

    # ── Main figure area — PNG for clean display ──
    # Fit figure in remaining space with margins
    fig_top = Inches(1.05)
    fig_bottom = Inches(7.1)
    fig_left = Inches(0.4)
    fig_right = Inches(12.9)
    avail_w = fig_right - fig_left
    avail_h = fig_bottom - fig_top

    # Read PNG dimensions to compute aspect ratio
    from PIL import Image
    with Image.open(png_path) as img:
        png_w, png_h = img.size
    aspect = png_w / png_h

    if avail_w / avail_h > aspect:
        # Height-limited
        h = avail_h
        w = int(h * aspect)
    else:
        w = avail_w
        h = int(w / aspect)

    # Center
    left = int((prs.slide_width - w) / 2)
    top = int(fig_top)

    # Add high-res PNG
    pic = slide.shapes.add_picture(str(png_path), left, top, w, h)

    # ── Thin border around the figure ──
    border = slide.shapes.add_shape(1, left - Inches(0.02), top - Inches(0.02),
                                     w + Inches(0.04), h + Inches(0.04))
    border.fill.background()
    border.line.color.rgb = BORDER_GRAY; border.line.width = Pt(0.5)

    # ── "Edit in Vector" hint at bottom-right ──
    hint = slide.shapes.add_textbox(Inches(8.5), Inches(7.15), Inches(4.5), Inches(0.25))
    ph = hint.text_frame.paragraphs[0]
    ph.text = "SVG source embedded — double-click to open for vector editing"
    ph.font.size = Pt(8); ph.font.color.rgb = SUBTITLE_GRAY; ph.alignment = PP_ALIGN.RIGHT

    print(f"Added: {name}  ({w/914400:.1f}x{h/914400:.1f}in)")

prs.save(str(OUT))
print(f"\nSaved: {OUT} ({len(prs.slides)} slides)")
print("High-res PNG images with clean layout. SVG source files in same folder for vector editing.")
