#!/usr/bin/env python3
"""Generate Section 4: Trustworthiness Evaluation PPTX slides for CMM."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── color palette ────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT  = RGBColor(0x00, 0x7A, 0xCC)   # blue
RED     = RGBColor(0xB4, 0x28, 0x28)   # entropy red
GREEN   = RGBColor(0x28, 0x8C, 0x50)   # gain green
GRAY    = RGBColor(0x66, 0x66, 0x66)
LGRAY   = RGBColor(0xE8, 0xE8, 0xE8)
YELLOW  = RGBColor(0xF0, 0xBC, 0x20)
ORANGE  = RGBColor(0xE0, 0x7A, 0x20)
CODE_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── helpers ──────────────────────────────────────────────────

def _blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def _bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _textbox(slide, left, top, width, height):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    return shape.text_frame

def _rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.line.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape

def _set_text(frame, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font='Calibri'):
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align

def _add_para(frame, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=4, space_after=2):
    p = frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    return p

def _title_slide(slide, title, subtitle=None):
    _rect(slide, 0, 0, 13.333, 0.08, ACCENT)
    tb = _textbox(slide, 0.7, 0.35, 12, 0.7)
    _set_text(tb, title, size=28, bold=True, color=WHITE)
    if subtitle:
        tb2 = _textbox(slide, 0.7, 1.05, 12, 0.45)
        _set_text(tb2, subtitle, size=14, color=GRAY)

def _section_title(slide, title):
    _rect(slide, 0, 0, 13.333, 0.06, ACCENT)
    tb = _textbox(slide, 0.7, 2.8, 11.9, 1.0)
    _set_text(tb, title, size=36, bold=True, color=WHITE)

def _code_block(slide, left, top, width, height, lines):
    shape = _rect(slide, left, top, width, height, fill=RGBColor(0x1E,0x1E,0x2E))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top  = Pt(6)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = 'Consolas'
        p.font.color.rgb = RGBColor(0x8E,0xC8,0x70)
        p.space_before = Pt(1)
        p.space_after  = Pt(1)

def _mini_table(slide, left, top, width, rows_data, col_widths=None, header=True):
    """rows_data: list of lists of strings; first row is header if header=True."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(0.35 * n_rows))
    table = table_shape.table
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
            if ri == 0 and header:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            else:
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2A,0x2A,0x3E)
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(cw)

# ═══════════════════════════════════════════════════════════════
# SLIDE 0 — Title
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_rect(s, 0, 0, 13.333, 0.10, ACCENT)
# decorative line
_rect(s, 0, 3.0, 13.333, 0.003, ACCENT)
tb = _textbox(s, 1.2, 1.4, 11.0, 1.2)
_set_text(tb, "Section 4: Trustworthiness Evaluation", size=40, bold=True, align=PP_ALIGN.LEFT)
_add_para(tb, "Preventing Mismatch Errors from GNSS and Map Faults", size=18, color=GRAY)
_add_para(tb, "Covariance Map Matching (CMM)  |  Chenzhang Ning  |  2026", size=13, color=ACCENT, space_before=16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — The Confidence Gap
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "1 — The Confidence Gap in Map Matching")

ls = [
    "Traditional HMM map matching outputs only a maximum-likelihood path.",
    "We get the optimal edge sequence, but no per-point confidence metric.",
    "This is insufficient for safety-critical applications."
]
for i, line in enumerate(ls):
    tb = _textbox(s, 0.8, 1.6 + i*0.65, 7.5, 0.5)
    _set_text(tb, f"•  {line}", size=16, color=WHITE)

# right panel: what we get vs what we need
rect_got = _rect(s, 8.6, 1.5, 2.1, 2.5, fill=RGBColor(0x2A,0x2A,0x3E))
tf = rect_got.text_frame; tf.word_wrap = True; tf.margin_left = Pt(8)
_set_text(tf, "What we get", size=15, bold=True, color=ACCENT)
_add_para(tf, "Optimal edges", size=12, color=WHITE)
_add_para(tf, "Cumulative prob.", size=12, color=WHITE)
_add_para(tf, "No confidence", size=12, color=RED)

rect_need = _rect(s, 10.9, 1.5, 2.1, 2.5, fill=RGBColor(0x2A,0x2A,0x3E))
tf2 = rect_need.text_frame; tf2.word_wrap = True; tf2.margin_left = Pt(8)
_set_text(tf2, "What we need", size=15, bold=True, color=GREEN)
_add_para(tf2, "Certainty score", size=12, color=WHITE)
_add_para(tf2, "Ambiguity loc.", size=12, color=WHITE)
_add_para(tf2, "Fault diagnosis", size=12, color=WHITE)

# bottom question
_rect(s, 0.8, 4.8, 11.5, 0.003, ACCENT)
tbq = _textbox(s, 0.8, 5.1, 11.5, 0.9)
_set_text(tbq, "Core Question: Can we quantify epistemic uncertainty in the matching decision at each epoch, using the same probabilistic framework that drove the match?", size=17, bold=True, color=YELLOW)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Two Complementary Signals
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "2 — Two Complementary Trustworthiness Signals")

# box 1
b1 = _rect(s, 1.0, 1.5, 11.3, 1.2, fill=RGBColor(0x0A,0x2A,0x4A))
tf = b1.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Layer Entropy & Delta Entropy", size=20, bold=True, color=ACCENT)
_add_para(tf, "H_post(t) — distributional uncertainty at epoch t    |    ΔH(t) — information gained from observation z_t", size=14, color=RGBColor(0xAA,0xCC,0xEE))
_add_para(tf, "• Instantaneous — measures local ambiguity at a single GPS point", size=13, color=WHITE, space_before=4)

# box 2
b2 = _rect(s, 1.0, 2.9, 11.3, 1.2, fill=RGBColor(0x4A,0x0A,0x0A))
tf = b2.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Sliding-Window Path Margin", size=20, bold=True, color=RED)
_add_para(tf, "margin(t) = p_rank1 − p_rank2   |   Top-3 beam-search over a contextual window", size=14, color=RGBColor(0xEE,0xAA,0xAA))
_add_para(tf, "• Contextual — detects path-level ambiguity that entropy alone misses", size=13, color=WHITE, space_before=4)

# box 3
b3 = _rect(s, 1.0, 4.3, 11.3, 1.2, fill=RGBColor(0x0A,0x3A,0x0A))
tf = b3.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Threshold Filtering", size=20, bold=True, color=GREEN)
_add_para(tf, "Discard epochs where margin exceeds threshold (likely GPS multipath or map digitization error)", size=14, color=RGBColor(0xAA,0xEE,0xAA))
_add_para(tf, "• Post-processing gate — enables safety-critical decision boundaries", size=13, color=WHITE, space_before=4)

# annotation
tb = _textbox(s, 1.0, 5.8, 11.0, 0.8)
_set_text(tb, "Together they distinguish local ambiguity from path-level confusion — orthogonal dimensions of uncertainty.", size=15, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Layer Entropy
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "3 — Layer Entropy as Uncertainty Measure")

# Formula block
tb = _textbox(s, 0.8, 1.5, 11.5, 0.5)
_set_text(tb, "For layer with N valid candidates having log-cumulative probabilities {ℓ₁, …, ℓₙ}:", size=16, color=WHITE)

# LogSumExp
b = _rect(s, 0.8, 2.1, 11.5, 0.85, fill=RGBColor(0x2A,0x2A,0x3E))
tf = b.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Step 1 — LogSumExp normalization (numerically stable):", size=14, bold=True, color=ACCENT)
_add_para(tf, "ℓ_max = max(ℓᵢ);   S = log Σ exp(ℓᵢ − ℓ_max);   log pᵢ_norm = ℓᵢ − ℓ_max − S", size=15, color=WHITE, space_before=6)

# Shannon
b2 = _rect(s, 0.8, 3.15, 11.5, 0.75, fill=RGBColor(0x2A,0x2A,0x3E))
tf = b2.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Step 2 — Shannon entropy (bits):    H = −Σ pᵢ_norm · log₂(pᵢ_norm)", size=16, bold=True, color=WHITE)
_add_para(tf, "H ≈ 0 → one candidate dominates (high confidence)   |   H ≈ log₂(N) → uniform distribution (maximum ambiguity)", size=14, color=GRAY, space_before=4)

# Code
_code_block(s, 0.8, 4.2, 11.5, 1.4, [
    "// Source: cmm_algorithm.cpp:1741-1763",
    "double log_sum = log_sum_exp(layer_log_probs);",
    "double inv_log2 = 1.0 / std::log(2.0);",
    "for (double log_p : layer_log_probs) {",
    "    double log_norm = log_p - log_sum;",
    "    double p_norm = std::exp(log_norm);",
    "    if (p_norm > 0.0)  layer_entropy -= p_norm * log_norm * inv_log2;",
    "}",
])

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Delta Entropy
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "4 — Information Gain: ΔH = H_prior − H_posterior")

tb = _textbox(s, 0.8, 1.5, 11.5, 0.4)
_set_text(tb, "Delta entropy quantifies how much the current GPS observation z_t reduces uncertainty.", size=16, color=GRAY)

# Left panel - first epoch
l = _rect(s, 0.8, 2.1, 5.6, 1.8, fill=RGBColor(0x2A,0x2A,0x3E))
tf = l.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "First Epoch (t = 0)", size=18, bold=True, color=ACCENT)
_add_para(tf, "No prior from propagation → uniform prior", size=14, color=WHITE, space_before=6)
_add_para(tf, "H_prior = log₂(N_valid)", size=14, color=WHITE)
_add_para(tf, "ΔH = log₂(N_valid) − H_posterior", size=14, color=GREEN, bold=True)

# Right panel - subsequent
r = _rect(s, 6.7, 2.1, 5.6, 1.8, fill=RGBColor(0x2A,0x2A,0x3E))
tf = r.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "Subsequent Epochs (t > 0)", size=18, bold=True, color=ACCENT)
_add_para(tf, "Prior = prediction from previous layer", size=14, color=WHITE, space_before=6)
_add_para(tf, "P_prior(b) ∝ Σ_a P(a)·TP(a,b)", size=14, color=WHITE)
_add_para(tf, "ΔH = H_prior − H_posterior", size=14, color=GREEN, bold=True)

# Big formula
fb = _rect(s, 1.5, 4.2, 10.3, 1.0, fill=RGBColor(0x0A,0x25,0x0A))
tf = fb.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "ΔH(t) = H[ P(x_t | z_{1:t-1}) ]  −  H[ P(x_t | z_{1:t}) ]", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_para(tf, "prediction entropy                      update (posterior) entropy", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# interpretation
tb2 = _textbox(s, 0.8, 5.5, 11.5, 1.0)
_set_text(tb2, "ΔH ≫ 0:  observation strongly disambiguates — informative   |   ΔH ≈ 0:  observation adds little — redundant or noisy", size=16, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — First Layer Code + Example
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "5 — Delta Entropy: First Layer (initialize_first_layer)")

_code_block(s, 0.8, 1.5, 7.2, 1.6, [
    "// Source: cmm_algorithm.cpp:1590-1625",
    "// Step 1: posterior entropy from cumu_probs",
    "(same Shannon entropy as before)",
    "",
    "// Step 2: uniform prior",
    "double H_prior = std::log2(static_cast<double>(",
    "    layer_log_probs.size()));",
    "delta_layer_entropy = H_prior - layer_entropy;",
    "if (delta_layer_entropy < 0.0)",
    "    delta_layer_entropy = 0.0;",
])

# example table on right
_rect(s, 8.3, 1.5, 4.5, 0.35, fill=ACCENT)
tb = _textbox(s, 8.3, 1.5, 4.5, 0.35)
_set_text(tb, "Example Scenarios", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_mini_table(s, 8.3, 1.85, 4.5, [
    ["Scenario",      "N", "ΔH (bits)"],
    ["Dominant (99%)", "8", "≈ 3.0"],
    ["Uniform over 4", "4", "≈ 0.0"],
    ["Mid (top 60%)",  "5", "≈ 1.2"],
], col_widths=[1.5, 0.8, 1.2])

# bottom interpretation
tb2 = _textbox(s, 0.8, 3.5, 11.5, 1.5)
_set_text(tb2, "Why ΔH ≥ 0 is clamped:", size=16, bold=True, color=YELLOW)
_add_para(tb2, "• Entropy cannot decrease below zero — a uniform prior has maximum entropy.", size=14, color=WHITE)
_add_para(tb2, "• ΔH < 0 would imply the posterior is more uncertain than the prior — numerical artifact.", size=14, color=WHITE)
_add_para(tb2, "• The first layer always starts from maximal uncertainty (uniform).", size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Subsequent Layers
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "6 — Delta Entropy: Subsequent Layers (update_layer_cmm)")

# diagram: prior → posterior
for i, (label, clr) in enumerate([
    ("Prediction\nP(x_t | z_{1:t-1})", ACCENT),
    ("× EP(z_t)", YELLOW),
    ("Posterior\nP(x_t | z_{1:t})", GREEN),
]):
    r = _rect(s, 1.0 + i*4, 1.5, 3.6, 1.0, fill=RGBColor(0x2A,0x2A,0x3E))
    tb = r.text_frame
    tb.word_wrap = True
    _set_text(tb, label, size=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
    if i < 2:
        # arrow between
        tb2 = _textbox(s, 4.6 + (i-1)*4, 1.8, 1.0, 0.5)
        _set_text(tb2, "→", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

_code_block(s, 0.8, 2.9, 11.5, 1.2, [
    "// During update_layer_cmm, accumulate prediction log-probs before multiplying by EP",
    "double log_sum_prev_probs = log_sum_exp(incoming_log_probs);",
    "layer_prior_log_probs.push_back(log_sum_prev_probs);",
    "",
    "// Then combine with EP to get posterior:",
    "node_b.cumu_prob = log_sum_prev_probs + node_b.ep;  // posterior",
    "",
    "// Separate entropy on layer_prior_log_probs → H_prior",
    "// Separate entropy on layer_log_probs → H_posterior",
    "// ΔH = H_prior − H_posterior  (cmm_algorithm.cpp:1781)",
])

# insight
tb3 = _textbox(s, 0.8, 4.5, 11.5, 1.2)
_set_text(tb3, "Key Insight:", size=18, bold=True, color=YELLOW)
_add_para(tb3, "ΔH isolates the marginal contribution of each epoch's GPS observation to the inference.", size=15, color=WHITE)
_add_para(tb3, "A sudden drop in ΔH signals degraded GNSS quality (e.g., entering an urban canyon).", size=15, color=RED)
_add_para(tb3, "• This enables real-time GNSS fault detection within the matching pipeline itself.", size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Why Window-Based?
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "7 — Why a Window-Based Metric?")

tb = _textbox(s, 0.8, 1.5, 11.5, 1.0)
_set_text(tb, "Single-epoch entropy is myopic:", size=18, bold=True, color=RED)
_add_para(tb, "•  A confident emission at t may lead to a wrong downstream path", size=15, color=WHITE)
_add_para(tb, "•  Ambiguity only becomes visible when considering the full sub-trajectory", size=15, color=WHITE)

# sliding window diagram
y_base = 3.0
for wi, (start, end, color, title, y) in enumerate([
    (1,  7,  ACCENT, "window (len=6), centered at t_7", 3.0),
    (3,  9,  RED,    "window (len=6), centered at t_9", 4.3),
    (5, 11,  GREEN,  "window (len=6), centered at t_11",5.6),
]):
    # dots
    for i in range(1, 13):
        label = f"t{i}"
        x = 1.0 + (i-1) * 0.95
        rect = _rect(s, x, y, 0.7, 0.5,
                      fill=color if start <= i <= end else RGBColor(0x2A,0x2A,0x3E))
        tf = rect.text_frame
        tf.word_wrap = True
        _set_text(tf, label, size=11, bold=True, color=WHITE if start<=i<=end else GRAY, align=PP_ALIGN.CENTER)
    # title
    tb = _textbox(s, 1.0, y - 0.35, 4.0, 0.3)
    _set_text(tb, title, size=11, color=color)

# margin formula
fb = _rect(s, 1.5, 6.4, 10.3, 0.6, fill=RGBColor(0x0A,0x25,0x0A))
tf = fb.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "For each window [t−k+1, t]:  beam-search top-3 path scores → margin(t) = p_rank1 − p_rank2", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Window Algorithm
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "8 — Sliding-Window Algorithm (compute_window_trustworthiness)")

_code_block(s, 0.5, 1.5, 7.8, 5.0, [
    "// Source: cmm_algorithm.cpp:1801-1900      (top-K = 3)",
    "",
    "for (size_t end_idx = 0; end_idx < layer_count; ++end_idx) {",
    "    start_idx = max(0, end_idx + 1 - window_length);",
    "",
    "    // 1. Seed: emission log-probs at start_idx",
    "    for (each candidate c at start) push_top_k(&scores[c], log_ep, 3);",
    "",
    "    // 2. Forward beam through window [start+1...end]",
    "    for (cursor = start+1; cursor <= end_idx; ++cursor) {",
    "        for (candidate b at cursor) {",
    "            for (candidate a at cursor-1) {",
    "                sp = get_sp_dist(a, b);",
    "                if (sp < 0) continue;",
    "                tp = calc_tp(sp, euclidean_distance[cursor-1]);",
    "                for each (prev_log in prev_scores[a]) {",
    "                    push_top_k(&cur_scores[b], prev_log+log(tp)+log_ep[b], 3);",
    "}   }   }   }",
    "    // 3. Collect combined top-3 from all ending candidates",
    "    // 4. LogSumExp normalize → linear probs → margin = combined[0]-combined[1]",
    "}",
])

# complexity box
r = _rect(s, 8.8, 1.5, 4.0, 2.5, fill=RGBColor(0x2A,0x2A,0x3E))
tf = r.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "Complexity Analysis", size=18, bold=True, color=YELLOW)
_add_para(tf, "", size=4)
_add_para(tf, "Exhaustive: O(K^W)", size=16, color=RED, bold=True, space_before=6)
_add_para(tf, "  K=8, W=10 → 8^10 ≈ 1e9", size=13, color=GRAY)
_add_para(tf, "", size=4, space_before=6)
_add_para(tf, "Beam (k=3): O(W·K²·k)", size=16, color=GREEN, bold=True, space_before=6)
_add_para(tf, "  10·64·3 = 1920 ops", size=13, color=GRAY)
_add_para(tf, "", size=4, space_before=6)
_add_para(tf, "> Online-capable", size=14, color=WHITE, space_before=6)

# bottom annotation
tb = _textbox(s, 0.8, 6.7, 11.5, 0.4)
_set_text(tb, "k=3 balances path diversity vs. computation   |   k=1 → no margin   |   k=2 → margin but may miss a 3rd distinct path   |   k>3 → diminishing returns", size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Margin in Action + Interpretation
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "9 — Trustworthiness Margin in Action")

_mini_table(s, 1.5, 1.6, 10.3, [
    ["Epoch", "rank-1 prob", "rank-2 prob", "Margin", "Interpretation", "Action"],
    ["t₁", "0.95", "0.03", "0.92", "Strong confidence", "Trust"],
    ["t₂", "0.48", "0.40", "0.08", "Near-tie", "Flag for review"],
    ["t₃", "0.55", "0.44", "0.11", "Moderate ambiguity", "Flag for review"],
    ["t₄", "0.88", "0.09", "0.79", "Strong confidence", "Trust"],
    ["t₅", "0.35", "0.33", "0.02", "High ambiguity", "Discard"],
], col_widths=[0.8, 1.2, 1.2, 0.9, 2.0, 1.4])

# guidelines
tb = _textbox(s, 0.8, 4.0, 11.5, 1.5)
_set_text(tb, "Interpretation Guidelines:", size=18, bold=True, color=YELLOW)
_add_para(tb, "•  margin > 0.5:  strong evidence for best path — can safely trust the match", size=15, color=GREEN)
_add_para(tb, "•  margin ∈ [0.1, 0.5]:  moderate — review context before acting", size=15, color=ORANGE)
_add_para(tb, "•  margin < 0.1:  near-tie — likely GNSS multipath or map parallel road ambiguity", size=15, color=RED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Threshold Filtering Pipeline
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "10 — Trustworthiness Threshold Filtering Pipeline")

# pipeline boxes
steps = [
    ("Viterbi\nbacktrack", ACCENT),
    ("Window\ntrustworthiness", YELLOW),
    ("Threshold\ncheck", RED),
    ("Output\nfiltered result", GREEN),
]

tb = _textbox(s, 0.8, 1.5, 11.5, 0.4)
_set_text(tb, "After Viterbi backtracking and window trustworthiness computation, a post-processing filter removes low-confidence epochs:", size=16, color=GRAY)

for i, (label, clr) in enumerate(steps):
    x = 1.0 + i * 3.0
    r = _rect(s, x, 2.3, 2.5, 1.0, fill=RGBColor(0x2A,0x2A,0x3E))
    tf = r.text_frame; tf.word_wrap = True
    _set_text(tf, label, size=15, bold=True, color=clr, align=PP_ALIGN.CENTER)
    if i < 3:
        tb2 = _textbox(s, x + 2.5, 2.55, 0.5, 0.4)
        _set_text(tb2, "→", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# code
_code_block(s, 0.8, 3.8, 11.5, 1.2, [
    "// Source: cmm_algorithm.cpp:1315-1323",
    "if (!config.filtered || trust <= config.trustworthiness_threshold) {",
    "    filtered_path.push_back(mc);               // keep epoch",
    "    filtered_tg_opath.push_back(node);          // keep node for path",
    "    filtered_indices.push_back(sub_indices[i]); // track original index",
    "}",
])

tb3 = _textbox(s, 0.8, 5.3, 11.5, 1.0)
_set_text(tb3, "Output fields preserved:", size=16, bold=True, color=YELLOW)
_add_para(tb3, "•  original_indices: pre-filter epoch indices for traceability", size=14, color=WHITE)
_add_para(tb3, "•  nbest_trustworthiness: top-N log scores per point for post-hoc analysis", size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Output Data Model
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "11 — Output Data Model (MatchedCandidate)")

# struct display
_code_block(s, 0.8, 1.5, 6.5, 2.3, [
    "// mm_type.hpp:64-72",
    "struct MatchedCandidate {",
    "    Candidate c;               // candidate edge projection",
    "    double ep;                 // emission probability",
    "    double tp;                 // transition probability",
    "    double cumu_prob;          // cumulative log-prob",
    "    double sp_dist;            // shortest-path distance",
    "    double trustworthiness;    // window margin (CMM)",
    "    double delta_entropy;      // ΔH in bits (CMM)",
    "};",
])

# field descriptions
_rect(s, 7.8, 1.5, 5.0, 0.35, fill=ACCENT)
tb = _textbox(s, 7.8, 1.55, 5.0, 0.3)
_set_text(tb, "Per-Epoch Fields", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_mini_table(s, 7.8, 1.85, 5.0, [
    ["Field", "Meaning"],
    ["trustworthiness", "Window margin (p1−p2)"],
    ["delta_entropy", "Info gain in bits"],
    ["cumu_prob", "Cumulative log-prob"],
    ["ep", "Emission probability"],
    ["tp", "Normalized transition"],
    ["sp_dist", "SP distance to prev"],
], col_widths=[1.8, 2.0])

# aggregate
_rect(s, 7.8, 4.2, 5.0, 0.35, fill=ACCENT)
tb = _textbox(s, 7.8, 4.25, 5.0, 0.3)
_set_text(tb, "Per-Result Aggregates", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_mini_table(s, 7.8, 4.55, 5.0, [
    ["Field", "Meaning"],
    ["nbest_trustworthiness", "Top-N scores/point"],
    ["original_indices", "Pre-filter indices"],
    ["candidate_details", "All candidates + EP"],
    ["status", "SUCCESS/PARTIAL/FAIL"],
], col_widths=[1.8, 2.0])

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Diagnostic Use Cases
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "12 — Diagnostic Use Cases: Fault Classification")

# left case
r = _rect(s, 0.8, 1.5, 5.8, 2.0, fill=RGBColor(0x2A,0x2A,0x3E))
tf = r.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "Case 1: Urban Canyon Multipath", size=18, bold=True, color=RED)
_add_para(tf, "• Dense high-rises → severe multipath", size=14, color=WHITE, space_before=6)
_add_para(tf, "• GPS errors become highly anisotropic", size=14, color=WHITE)
_add_para(tf, "• Cov ellipse covers multiple parallel roads → low margin", size=14, color=WHITE)
_add_para(tf, "• ΔH drops sharply vs. open-sky segments", size=14, color=YELLOW)

r2 = _rect(s, 6.9, 1.5, 5.8, 2.0, fill=RGBColor(0x2A,0x2A,0x3E))
tf = r2.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10)
_set_text(tf, "Case 2: Map Digitization Error", size=18, bold=True, color=ORANGE)
_add_para(tf, "• Road centerline offset from true trajectory", size=14, color=WHITE, space_before=6)
_add_para(tf, "• Consistent GPS but geometric mismatch", size=14, color=WHITE)
_add_para(tf, "• Multiple nearby edges have similar EP → persistently low margin", size=14, color=WHITE)
_add_para(tf, "• ΔH is normal (no GPS anomaly)", size=14, color=YELLOW)

# fault classification matrix
_mini_table(s, 0.8, 3.9, 11.5, [
    ["Fault Type",          "Margin",              "ΔH",             "Pattern",         "Root Cause"],
    ["GPS Multipath",       "Low (isolated)",      "Sharp drop",    "Spike",           "GNSS"],
    ["GPS Signal Loss",     "Zero (multiple)",     "Zeros",         "Gap segment",     "GNSS"],
    ["Map Offset",          "Low (persistent)",    "Normal",        "Sustained",       "Map"],
    ["Map Connectivity",    "Low (local)+disconn", "Normal",        "Junction",        "Map"],
    ["Healthy",             "High (>0.5)",         "High (>1.0)",   "Stable",          "—"],
], col_widths=[1.8, 1.8, 1.6, 1.5, 1.2])

tb = _textbox(s, 0.8, 6.2, 11.5, 0.8)
_set_text(tb, "Combinational analysis of margin + ΔH allows fault-type diagnosis: distinguishes GPS faults from map faults.", size=16, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Summary
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_title_slide(s, "13 — Summary: Trustworthiness in CMM")

items = [
    ("1", "Two complementary metrics capture different uncertainty dimensions:",
     "Layer entropy/delta-entropy (instantaneous) + Sliding-window path margin (contextual)"),
    ("2", "Computationally tractable:",
     "Beam-search (k=3) prevents exponential blowup — O(W·K²·k) per window"),
    ("3", "Diagnostically rich:",
     "Combinational analysis of margin + ΔH classifies faults (GPS vs. map vs. both)"),
    ("4", "Safety enabling:",
     "Trustworthiness thresholds let downstream systems decide confidence before acting"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.5 + i*1.25
    r = _rect(s, 0.8, y, 0.6, 0.6, fill=ACCENT)
    tf = r.text_frame
    _set_text(tf, num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb = _textbox(s, 1.7, y, 11.0, 0.4)
    _set_text(tb, title, size=18, bold=True, color=WHITE)
    tb2 = _textbox(s, 1.7, y+0.4, 11.0, 0.5)
    _set_text(tb2, desc, size=14, color=GRAY)

# takeaway
r = _rect(s, 1.0, 6.2, 11.3, 0.7, fill=RGBColor(0x0A,0x15,0x3A))
tf = r.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12)
_set_text(tf, "Takeaway: CMM's trustworthiness evaluation transforms map matching from a black-box decision into a quantified inference with per-epoch confidence bounds — essential for safety-critical applications.", size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════
s = _blank_slide(); _bg(s)
_rect(s, 0, 0, 13.333, 0.08, ACCENT)
tb = _textbox(s, 0.7, 2.6, 12, 1.5)
_set_text(tb, "Thank You", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb2 = _textbox(s, 0.7, 4.2, 12, 0.8)
_set_text(tb2, "Questions & Discussion", size=24, color=GRAY, align=PP_ALIGN.CENTER)

# ── save ──────────────────────────────────────────────────────
outpath = "docs/cmm_section4_trustworthiness.pptx"
prs.save(outpath)
print(f"Saved {len(prs.slides)} slides to {outpath}")
